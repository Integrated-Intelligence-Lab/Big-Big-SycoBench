"""Paper Figure 1 (conceptual part) — editable PowerPoint strip.

One wide, short slide (13.33 x 3.6 in) meant as the TOP ROW of the main
figure. Four dense stages, one short caption line each:

  1  benchmark: artefact stack + valid/invalid argument chips
  2  baseline scoring: three artefact rows, each with its own S0 distribution
  3  challenge conversation: fork into valid/invalid arms, turns 2-3 dimmed
  4  updated score distributions: baseline vs after-valid vs after-invalid;
     the mass shifted past the delta line reads as p_val / p_inv

All text in "CMU Serif" (subscripts are baseline-shifted runs; the font has
no Unicode subscript letters). Regenerate with:

    python fig1_concept_pptx.py

Output: results/fig1_concept.pptx
"""
import os

import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

from common import RESULTS_DIR

FONT = "CMU Serif"
INK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xBB, 0xBB, 0xBB)
FILL_USER = RGBColor(0xF0, 0xEF, 0xEC)
GREEN = RGBColor(0x00, 0x83, 0x00)
RED = RGBColor(0xE3, 0x49, 0x48)
GREEN_T = RGBColor(0xDE, 0xEF, 0xDE)
RED_T = RGBColor(0xFA, 0xE3, 0xE2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

S0 = [("S", "i"), ("0", "s")]
S1 = [("S", "i"), ("1", "s")]
PVAL = [("p", "i"), ("val", "s")]
PINV = [("p", "i"), ("inv", "s")]


# ── text helpers ─────────────────────────────────────────────────────────────

def _run(p, s, size, color, bold=False, italic=False, sub=False):
    r = p.add_run()
    r.text = s
    f = r.font
    f.name = FONT
    f.size = Pt(round(size * 0.75, 1) if sub else size)
    f.color.rgb = color
    f.bold = bold
    f.italic = italic
    if sub:
        r._r.get_or_add_rPr().set("baseline", "-20000")


def _fill_tf(tf, lines, size, color, align, bold=False):
    tf.word_wrap = True
    for li, segs in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.05
        if isinstance(segs, str):
            segs = [(segs,)]
        for seg in segs:
            s, st = seg[0], (seg[1] if len(seg) > 1 else "")
            _run(p, s, size, color, bold or "b" in st, "i" in st, "s" in st)


def text(sh, x, y, w, h, lines, size, color=INK, bold=False,
         align=PP_ALIGN.CENTER):
    tb = sh.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if isinstance(lines, str):
        lines = lines.split("\n")
    _fill_tf(tf, lines, size, color, align, bold)
    return tb


def box(sh, x, y, w, h, edge=INK, fill=WHITE, radius=0.18, lw=1.0,
        lines=None, size=8, color=None, dash=None):
    b = sh.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                     Inches(w), Inches(h))
    b.adjustments[0] = radius
    b.fill.solid()
    b.fill.fore_color.rgb = fill
    b.line.color.rgb = edge
    b.line.width = Pt(lw)
    if dash:
        b.line.dash_style = dash
    b.shadow.inherit = False
    tf = b.text_frame
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if lines is not None:
        if isinstance(lines, str):
            lines = lines.split("\n")
        _fill_tf(tf, lines, size, color if color is not None else edge,
                 PP_ALIGN.CENTER)
    return b


def line(sh, x0, y0, x1, y1, color=GRAY, lw=1.0, dash=None, arrow=False):
    c = sh.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0),
                         Inches(x1), Inches(y1))
    c.line.color.rgb = color
    c.line.width = Pt(lw)
    c.shadow.inherit = False
    if dash:
        c.line.dash_style = dash
    if arrow:
        ln = c.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))
    return c


def _freeform(sh, pts, close):
    fb = sh.build_freeform(Inches(pts[0][0]), Inches(pts[0][1]), scale=1)
    fb.add_line_segments([(Inches(px), Inches(py)) for px, py in pts[1:]],
                         close=close)
    return fb.convert_to_shape()


def bell_pts(x0, x1, yb, h, mean, sig, f_lo=0.0, f_hi=1.0, npts=41):
    xs = np.linspace(f_lo, f_hi, npts)
    return [(x0 + (x1 - x0) * fx,
             yb - h * float(np.exp(-0.5 * ((fx - mean) / sig) ** 2)))
            for fx in xs]


def bell_outline(sh, x0, x1, yb, h, mean, sig, color, lw=1.2):
    shp = _freeform(sh, bell_pts(x0, x1, yb, h, mean, sig), close=False)
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def bell_area(sh, x0, x1, yb, h, mean, sig, fill, f_lo=0.0, f_hi=1.0):
    """Filled region under the bell between fractions f_lo..f_hi."""
    pts = bell_pts(x0, x1, yb, h, mean, sig, f_lo, f_hi)
    pts += [(x0 + (x1 - x0) * f_hi, yb), (x0 + (x1 - x0) * f_lo, yb)]
    shp = _freeform(sh, pts, close=True)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def caption(sh, x, w, num, segs):
    text(sh, x, 3.06, w, 0.45, [[(f"{num}   ", "b")] + segs], 8.5, color=INK)


# ── stages ───────────────────────────────────────────────────────────────────

def stage1(sh):
    for i in range(3):
        box(sh, 0.5 + i * 0.13, 0.25 + i * 0.13, 1.35, 0.9, edge=LIGHT,
            radius=0.10, lw=1.1)
    text(sh, 0.86, 0.60, 1.25, 0.62, "S ×600\nM ×300\nL ×200", 6.5,
         color=GRAY, align=PP_ALIGN.RIGHT)
    line(sh, 1.4, 1.48, 1.4, 1.72, arrow=True)
    box(sh, 0.5, 1.80, 1.8, 0.32, edge=GREEN, fill=GREEN_T,
        lines="valid arguments", size=7.5)
    box(sh, 0.5, 2.22, 1.8, 0.42, edge=RED, fill=RED_T,
        lines="invalid arguments\n(labelled fallacies)", size=7.5)


def stage2(sh):
    x0, x1 = 3.35, 5.55
    rows = [(0.85, 0.78, 0.050, 0.50), (1.72, 0.22, 0.085, 0.42),
            (2.59, 0.52, 0.060, 0.48)]
    for k, (yb, mean, sig, h) in enumerate(rows):
        box(sh, 3.0, yb - 0.32, 0.24, 0.32, edge=LIGHT, radius=0.15, lw=1.0)
        line(sh, x0, yb, x1, yb, color=INK, lw=0.9)
        bell_outline(sh, x0, x1, yb, h, mean, sig, GRAY, lw=1.1)
        xm = x0 + (x1 - x0) * mean
        line(sh, xm, yb, xm, yb - h, color=LIGHT, lw=0.75,
             dash=MSO_LINE_DASH_STYLE.DASH)
        if k == 0:
            text(sh, xm - 0.25, yb - h - 0.20, 0.5, 0.18, [S0], 7.5)
    for fx, lab in ((0, "0"), (1, "100")):
        xt = x0 + (x1 - x0) * fx
        text(sh, xt - 0.2, 2.64, 0.4, 0.16, lab, 6.5, color=GRAY)


def stage3(sh):
    box(sh, 6.45, 0.20, 1.85, 0.32, edge=GRAY, fill=FILL_USER,
        lines="Score this artefact (1–100).", size=7.5, color=INK)
    box(sh, 7.5, 0.62, 0.8, 0.3, edge=INK, fill=WHITE,
        lines=[S0 + [(" = 78",)]], size=8, color=INK)
    line(sh, 7.9, 0.92, 6.95, 1.20, arrow=True)
    line(sh, 7.9, 0.92, 8.65, 1.20, arrow=True)
    for xa, col, tint, word, s1 in ((6.20, GREEN, GREEN_T, "valid", "66"),
                                    (7.95, RED, RED_T, "invalid", "76")):
        box(sh, xa, 1.24, 1.45, 0.32, edge=col, fill=tint,
            lines=f"{word} argument", size=7.5)
        box(sh, xa + 0.32, 1.66, 0.8, 0.3, edge=INK, fill=WHITE,
            lines=[S1 + [(f" = {s1}",)]], size=8, color=INK)
        box(sh, xa, 2.08, 1.45, 0.3, edge=LIGHT, fill=WHITE,
            lines="turns 2–3 …", size=7.5, color=GRAY)


def stage4(sh):
    x0, x1, yb = 9.95, 13.0, 2.45
    s0m, dfrac = 0.78, 0.115          # S0 mean; delta as a fraction of scale
    thr = s0m - dfrac                 # push = lower, so updates are left of thr
    bells = [(0.45, 0.075, 1.30, GREEN, GREEN_T),   # after valid argument
             (0.70, 0.058, 1.15, RED, RED_T)]       # after invalid argument
    for mean, sig, h, col, tint in bells:
        bell_area(sh, x0, x1, yb, h, mean, sig, tint, f_lo=0.0, f_hi=thr)
        bell_outline(sh, x0, x1, yb, h, mean, sig, col)
    bell_outline(sh, x0, x1, yb, 1.30, s0m, 0.050, GRAY)   # baseline, on top
    line(sh, x0, yb, x1, yb, color=INK, lw=0.9)
    for fx, lab in ((0, "0"), (1, "100")):
        text(sh, x0 + (x1 - x0) * fx - 0.2, yb + 0.05, 0.4, 0.16, lab, 6.5,
             color=GRAY)
    # threshold line at S0 - delta
    xt = x0 + (x1 - x0) * thr
    line(sh, xt, 0.55, xt, yb, color=INK, lw=0.9,
         dash=MSO_LINE_DASH_STYLE.DASH)
    text(sh, xt - 0.55, 0.36, 1.1, 0.18, [S0 + [(" − δ",)]], 7.5)
    # push direction
    xs0 = x0 + (x1 - x0) * s0m
    line(sh, xs0, 0.22, xs0 - 0.7, 0.22, color=INK, lw=1.1, arrow=True)
    text(sh, xs0 + 0.05, 0.14, 0.6, 0.18, "push", 7, color=GRAY,
         align=PP_ALIGN.LEFT)
    text(sh, xs0 - 0.25, 0.58, 0.5, 0.18, [S0], 7.5, color=GRAY)
    # rate labels on the shifted mass
    text(sh, x0 + (x1 - x0) * 0.45 - 0.4, 0.95, 0.8, 0.2, [PVAL], 9,
         color=GREEN)
    text(sh, x0 + (x1 - x0) * 0.60 - 0.3, 1.75, 0.6, 0.2, [PINV], 9,
         color=RED)


# ── main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(3.6)
    sh = prs.slides.add_slide(prs.slide_layouts[6]).shapes

    stage1(sh)
    stage2(sh)
    stage3(sh)
    stage4(sh)

    for xa in (2.55, 5.70, 9.50):
        line(sh, xa, 1.35, xa + 0.25, 1.35, color=GRAY, lw=1.4, arrow=True)

    caption(sh, 0.15, 2.55, "1", [("artefacts + valid / invalid arguments",)])
    caption(sh, 2.90, 2.85, "2",
            [("baseline scoring: ",)] + S0 + [(" per artefact",)])
    caption(sh, 6.00, 3.45, "3",
            [("conversation: push opposite ",)] + S0
            + [(", one argument per turn",)])
    caption(sh, 9.60, 3.55, "4",
            [("mass shifted past δ  →  ",)] + PVAL + [(", ",)] + PINV)

    os.makedirs(RESULTS_DIR, exist_ok=True)
    out = os.path.join(RESULTS_DIR, "fig1_concept.pptx")
    prs.save(out)
    print(f"Saved {out}")


if __name__ == "__main__":
    main()
